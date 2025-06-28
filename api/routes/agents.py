from enum import Enum
from logging import getLogger
from typing import AsyncGenerator, List, Optional, Union, Any, Dict
from datetime import datetime
import tempfile
from pathlib import Path
import base64
import os
import json
import asyncio

from agno.agent import Agent, AgentKnowledge
from agno.storage.session.agent import AgentSession
from agno.media import File, Image, Audio, Video
from fastapi import APIRouter, HTTPException, status, UploadFile, File as FastAPIFile, Form, Request, Depends
from fastapi.responses import StreamingResponse
from pydantic import BaseModel, Field

# Проверяем доступность библиотек для документов
try:
    from openpyxl import load_workbook
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    from pptx import Presentation
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False

try:
    from agno.document.reader.docx_reader import DocxReader
    AGNO_DOCX_AVAILABLE = True
except ImportError:
    AGNO_DOCX_AVAILABLE = False

try:
    from agno.document.reader.pdf_reader import PDFReader
    AGNO_PDF_AVAILABLE = True
except ImportError:
    AGNO_PDF_AVAILABLE = False

from agents.static.agno_assist import get_agno_assist_knowledge
from agents.selector import get_agent, get_available_agents, get_static_agent_details, is_static_agent, get_static_agents
from agents.models.saas_models import StaticAgentResponse
from db.session import SessionLocal
from sqlalchemy import text
from agents.cache import cache_manager
from openai import AsyncOpenAI

logger = getLogger(__name__)

# Константы для обработки документов
MAX_DOCUMENT_SIZE = 50 * 1024 * 1024  # 50MB максимальный размер документа
MAX_EXCEL_ROWS = 10000  # Максимальное количество строк для Excel
MAX_PPTX_SLIDES = 100  # Максимальное количество слайдов для PowerPoint

######################################################
## Routes for the Agent Interface
######################################################

agents_router = APIRouter(prefix="/agents", tags=["Agents"])


class Model(str, Enum):
    gpt_4_1 = "gpt-4.1"
    o4_mini = "o4-mini"


@agents_router.get("", response_model=List[Union[StaticAgentResponse, dict]])
async def list_agents():
    """
    Returns a list of all available agents with full information.
    Static agents return StaticAgentResponse format, dynamic agents return dict format.
    OPTIMIZED: Uses bulk operations to avoid N+1 queries, but respects cache for auto-updates.

    Returns:
        List[Union[StaticAgentResponse, dict]]: List of agent information
    """
    from agents.cache import cache_manager
    
    # Проверяем кэш для полного списка агентов
    cache_key = "agents:full_list"
    cached_agents = cache_manager.get(cache_key)
    if cached_agents:
        logger.info(f"Returned {len(cached_agents)} agents from cache")
        return cached_agents
    
    agents_info = []
    
    try:
        # 1. Получаем все статические агенты (быстрая операция)
        static_agent_ids = get_static_agents()
        
        # 2. Получаем всех динамических агентов одним запросом (вместо N запросов)
        dynamic_agents_bulk = []
        try:
            with SessionLocal() as session:
                query = text("""
                    SELECT id, name, agent_id, description, instructions,
                           model_config, tools_config, knowledge_config,
                           memory_config, storage_config, settings, is_active, 
                           created_at, updated_at
                    FROM dynamic_agents
                    WHERE is_active = true
                    ORDER BY created_at DESC
                """)
                
                result = session.execute(query)
                for row in result.fetchall():
                    # Извлекаем model_id из model_config
                    model_config = row.model_config if row.model_config else {}
                    model_id = model_config.get("id", "gpt-4o-mini")
                    
                    dynamic_agent_info = {
                        "id": row.id,
                        "name": row.name,
                        "agent_id": row.agent_id,
                        "description": row.description,
                        "instructions": row.instructions,
                        "model_id": model_id,
                        "model_config": model_config,
                        "tools_config": row.tools_config if row.tools_config else [],
                        "knowledge_config": row.knowledge_config if row.knowledge_config else {},
                        "memory_config": row.memory_config if row.memory_config else {},
                        "storage_config": row.storage_config if row.storage_config else {},
                        "settings": row.settings if row.settings else {},
                        "is_active": row.is_active,
                        "max_tokens": None,
                        "temperature": None,
                        "created_at": row.created_at.isoformat() if row.created_at else None,
                        "updated_at": row.updated_at.isoformat() if row.updated_at else None
                    }
                    dynamic_agents_bulk.append(dynamic_agent_info)
                    
        except Exception as db_error:
            logger.error(f"Database error while loading dynamic agents: {db_error}")
        
        # 3. Добавляем статические агенты (ОПТИМИЗИРОВАНО: без создания Agent объектов)
        for agent_id in static_agent_ids:
            try:
                # Быстрое получение базовой информации без создания Agent объекта
                from agents.selector import get_static_agent_basic_info
                basic_info = get_static_agent_basic_info(agent_id)
                
                if basic_info:
                    static_agent = StaticAgentResponse(
                        id=None,
                        name=basic_info.get("name", agent_id),
                        agent_id=agent_id,
                        description=basic_info.get("description"),
                        instructions=basic_info.get("instructions"),
                        model_id=basic_info.get("model_id", "gpt-4.1"),
                        model_config_data={"type": "openai", "id": basic_info.get("model_id", "gpt-4.1")},
                        tools_config=basic_info.get("tools_config", []),
                        knowledge_config=basic_info.get("knowledge_config", {}),
                        memory_config=basic_info.get("memory_config", {}),
                        storage_config=basic_info.get("storage_config", {}),
                        settings={},  # Базовые настройки для списка
                        is_active=basic_info.get("is_active", True),
                        max_tokens=None,
                        temperature=None,
                        created_at=None,
                        updated_at=None,
                        agent_type="static",
                        source_file=basic_info.get("source_file"),
                        editable=False
                    )
                    agents_info.append(static_agent.model_dump(by_alias=True))
                else:
                    # Fallback если не удалось получить информацию
                    agents_info.append({
                        "agent_id": agent_id,
                        "agent_type": "static",
                        "name": agent_id.replace('_', ' ').title(),
                        "is_active": True,
                        "editable": False
                    })
            except Exception as e:
                logger.warning(f"Failed to get basic info for static agent {agent_id}: {e}")
                # Добавляем минимальную информацию при ошибке
                agents_info.append({
                    "agent_id": agent_id,
                    "agent_type": "static",
                    "name": agent_id,
                    "error": str(e)
                })
        
        # 4. Добавляем все динамические агенты (уже загружены bulk операцией)
        agents_info.extend(dynamic_agents_bulk)
        
        # 5. Кэшируем результат на 5 минут (с поддержкой автообновления)
        cache_manager.set(cache_key, agents_info, ttl=300)
        
        logger.info(f"Loaded {len(static_agent_ids)} static and {len(dynamic_agents_bulk)} dynamic agents")
        
    except Exception as e:
        logger.error(f"Error in list_agents: {e}")
        # В случае критической ошибки возвращаем хотя бы базовый список
        return [{
            "error": "Failed to load agents",
            "message": str(e)
        }]
    
    return agents_info


async def chat_response_streamer(
    agent: Agent, 
    message: str, 
    files: Optional[List[File]] = None,
    images: Optional[List[Image]] = None,
    audio: Optional[List[Audio]] = None,
    videos: Optional[List[Video]] = None
) -> AsyncGenerator:
    """
    Stream agent responses chunk by chunk with full multimodal support.

    Args:
        agent: The agent instance to interact with
        message: User message to process
        files: Optional list of files to process
        images: Optional list of images to process
        audio: Optional list of audio to process
        videos: Optional list of videos to process

    Yields:
        JSON-formatted chunks with text content and multimedia artifacts
    """
    import json
    
    try:
        run_response = await agent.arun(
            message, 
            stream=True,
            files=files,
            images=images,
            audio=audio,
            videos=videos
        )
        
        # Отслеживаем накопленные мультимедиа артефакты
        accumulated_images = []
        accumulated_audio = []
        accumulated_videos = []
        response_audio = None
        
        async for chunk in run_response:
            try:
                # Базовая структура chunk ответа
                chunk_data = {
                    "type": "text",
                    "content": None
                }
                
                # Обрабатываем текстовое содержимое
                if chunk and hasattr(chunk, 'content') and chunk.content is not None:
                    chunk_data["content"] = chunk.content
                elif chunk and hasattr(chunk, 'delta') and chunk.delta:
                    chunk_data["content"] = str(chunk.delta)
                
                # 🆕 ПОДДЕРЖКА МУЛЬТИМЕДИА В STREAMING
                # Проверяем на новые изображения
                if hasattr(chunk, 'images') and chunk.images:
                    for img in chunk.images:
                        if img not in accumulated_images:
                            accumulated_images.append({
                                "id": getattr(img, 'id', None),
                                "url": getattr(img, 'url', None),
                                "content": getattr(img, 'content', None),
                                "mime_type": getattr(img, 'mime_type', None),
                                "alt_text": getattr(img, 'alt_text', None)
                            })
                            # Отправляем событие о новом изображении
                            yield json.dumps({
                                "type": "image",
                                "data": accumulated_images[-1]
                            }) + "\n"
                
                # Проверяем на новые аудио
                if hasattr(chunk, 'audio') and chunk.audio:
                    for aud in chunk.audio:
                        if aud not in accumulated_audio:
                            accumulated_audio.append({
                                "id": getattr(aud, 'id', None),
                                "url": getattr(aud, 'url', None),
                                "content": getattr(aud, 'base64_audio', None),
                                "mime_type": getattr(aud, 'mime_type', None),
                                "length": getattr(aud, 'length', None)
                            })
                            # Отправляем событие о новом аудио
                            yield json.dumps({
                                "type": "audio",
                                "data": accumulated_audio[-1]
                            }) + "\n"
                
                # Проверяем на новые видео
                if hasattr(chunk, 'videos') and chunk.videos:
                    for vid in chunk.videos:
                        if vid not in accumulated_videos:
                            accumulated_videos.append({
                                "id": getattr(vid, 'id', None),
                                "url": getattr(vid, 'url', None),
                                "content": getattr(vid, 'content', None),
                                "mime_type": getattr(vid, 'mime_type', None),
                                "eta": getattr(vid, 'eta', None),
                                "length": getattr(vid, 'length', None)
                            })
                            # Отправляем событие о новом видео
                            yield json.dumps({
                                "type": "video",
                                "data": accumulated_videos[-1]
                            }) + "\n"
                
                # Проверяем на аудио ответ
                if hasattr(chunk, 'response_audio') and chunk.response_audio and not response_audio:
                    response_audio = {
                        "id": getattr(chunk.response_audio, 'id', None),
                        "content": getattr(chunk.response_audio, 'content', None),
                        "expires_at": getattr(chunk.response_audio, 'expires_at', None),
                        "transcript": getattr(chunk.response_audio, 'transcript', None),
                        "mime_type": getattr(chunk.response_audio, 'mime_type', None),
                        "sample_rate": getattr(chunk.response_audio, 'sample_rate', None),
                        "channels": getattr(chunk.response_audio, 'channels', None)
                    }
                    # Отправляем событие об аудио ответе
                    yield json.dumps({
                        "type": "response_audio",
                        "data": response_audio
                    }) + "\n"
                
                # Отправляем текстовый chunk если есть содержимое
                if chunk_data["content"]:
                    yield json.dumps(chunk_data) + "\n"
                    
            except Exception as chunk_error:
                logger.warning(f"Error processing streaming chunk: {chunk_error}")
                continue
                
    except Exception as e:
        logger.error(f"Error in streaming response: {e}")
        yield json.dumps({
            "type": "error",
            "content": f"Ошибка при обработке запроса: {str(e)}"
        }) + "\n"


async def process_uploaded_files(
    files: Optional[List[UploadFile]] = None,
    images: Optional[List[UploadFile]] = None,
    audio: Optional[List[UploadFile]] = None,
    videos: Optional[List[UploadFile]] = None
) -> tuple[Optional[List[File]], Optional[List[Image]], Optional[List[Audio]], Optional[List[Video]], str]:
    """
    Обрабатывает загруженные файлы и конвертирует их в объекты Agno Media.
    Полная совместимость с Agno Agent arun() методом.

    Args:
        files: Загруженные файлы (документы, код, текст)
        images: Загруженные изображения 
        audio: Загруженные аудио файлы
        videos: Загруженные видео файлы

    Returns:
        Кортеж: (files, images, audio, videos, text_content_summary)
    """
    processed_files = []
    processed_images = []
    processed_audio = []
    processed_videos = []
    text_content_summary = ""

    # Поддерживаемые MIME типы для файлов в OpenAI (только PDF)
    OPENAI_SUPPORTED_FILE_TYPES = {"application/pdf"}

    def determine_file_mime_type(upload_file: UploadFile) -> str:
        """Определяет MIME тип файла на основе content_type и расширения"""
        content_type = upload_file.content_type
        filename = upload_file.filename or ""
        
        # Если MIME тип уже поддерживается OpenAI, используем его
        if content_type and content_type in OPENAI_SUPPORTED_FILE_TYPES:
            return content_type
            
        # Иначе определяем по расширению файла
        if filename and '.' in filename:
            ext = filename.lower().split('.')[-1]
            if ext == 'pdf':
                return 'application/pdf'
        
        # Для остальных типов возвращаем None (будем обрабатывать как текст)
        return None

    def is_text_file(filename: str, content_type: str) -> bool:
        """Проверяет, является ли файл текстовым"""
        if not filename:
            return False
        
        ext = filename.lower().split('.')[-1] if '.' in filename else ''
        text_extensions = {
            'txt', 'md', 'markdown', 'py', 'js', 'html', 'htm', 'css', 
            'json', 'xml', 'yaml', 'yml', 'csv', 'rtf', 'log'
        }
        
        return ext in text_extensions or (content_type and content_type.startswith('text/'))
    
    def is_document_file(filename: str) -> bool:
        """Проверяет, является ли файл документом (DOCX, XLSX, PPTX)"""
        if not filename:
            return False
        
        ext = filename.lower().split('.')[-1] if '.' in filename else ''
        document_extensions = {'docx', 'xlsx', 'pptx'}
        
        return ext in document_extensions
    
    async def process_document_file(content: bytes, filename: str) -> str:
        """Обрабатывает документы DOCX, XLSX, PPTX с улучшенной обработкой ошибок и ограничениями"""
        try:
            # Проверяем размер файла
            if len(content) > MAX_DOCUMENT_SIZE:
                return f"Документ слишком большой: {len(content):,} байт (максимум: {MAX_DOCUMENT_SIZE:,} байт)"
            
            ext = filename.lower().split('.')[-1] if '.' in filename else ''
            logger.info(f"Processing document: {filename} ({len(content):,} bytes, type: {ext})")
            
            if ext == 'docx':
                if not AGNO_DOCX_AVAILABLE:
                    return "Agno DocxReader недоступен. Проверьте установку библиотеки python-docx"
                
                reader = DocxReader()
                
                # Создаем временный файл для DocxReader
                import tempfile
                with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp_file:
                    try:
                        tmp_file.write(content)
                        tmp_file.flush()
                        
                        # Читаем документ через Agno
                        documents = reader.read(Path(tmp_file.name))
                        
                        if documents and len(documents) > 0:
                            doc_content = documents[0].content
                            if doc_content and len(doc_content.strip()) > 0:
                                return doc_content
                            else:
                                return "DOCX файл не содержит текстового содержимого"
                        else:
                            return "Не удалось извлечь содержимое из DOCX файла"
                    
                    finally:
                        # Всегда удаляем временный файл
                        try:
                            os.unlink(tmp_file.name)
                        except OSError:
                            pass
            
            elif ext == 'xlsx':
                if not OPENPYXL_AVAILABLE:
                    return "Библиотека openpyxl недоступна. Установите: pip install openpyxl"
                
                import io
                
                # Читаем Excel файл
                workbook = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
                
                try:
                    content_parts = []
                    total_rows = 0
                    
                    for sheet_name in workbook.sheetnames:
                        worksheet = workbook[sheet_name]
                        content_parts.append(f"Лист '{sheet_name}':")
                        
                        sheet_rows = 0
                        # Читаем данные построчно с ограничением
                        for row in worksheet.iter_rows(values_only=True):
                            if total_rows >= MAX_EXCEL_ROWS:
                                content_parts.append(f"... (достигнуто ограничение в {MAX_EXCEL_ROWS:,} строк)")
                                break
                            
                            # Фильтруем пустые строки
                            if any(cell is not None and str(cell).strip() for cell in row):
                                row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                                content_parts.append(row_text)
                                sheet_rows += 1
                                total_rows += 1
                        
                        content_parts.append(f"(Обработано строк: {sheet_rows})")
                        content_parts.append("")  # Пустая строка между листами
                        
                        if total_rows >= MAX_EXCEL_ROWS:
                            break
                    
                    return "\n".join(content_parts)
                
                finally:
                    workbook.close()
            
            elif ext == 'pptx':
                if not PYTHON_PPTX_AVAILABLE:
                    return "Библиотека python-pptx недоступна. Установите: pip install python-pptx"
                
                import io
                
                # Читаем презентацию
                prs = Presentation(io.BytesIO(content))
                
                content_parts = []
                slide_count = len(prs.slides)
                
                if slide_count > MAX_PPTX_SLIDES:
                    content_parts.append(f"Презентация содержит {slide_count} слайдов, обрабатываем первые {MAX_PPTX_SLIDES}")
                    slides_to_process = prs.slides[:MAX_PPTX_SLIDES]
                else:
                    slides_to_process = prs.slides
                
                for i, slide in enumerate(slides_to_process, 1):
                    content_parts.append(f"Слайд {i}:")
                    
                    slide_text_found = False
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            content_parts.append(f"  {shape.text.strip()}")
                            slide_text_found = True
                    
                    if not slide_text_found:
                        content_parts.append("  (Текстовое содержимое не найдено)")
                    
                    content_parts.append("")
                
                content_parts.append(f"Обработано слайдов: {len(slides_to_process)} из {slide_count}")
                return "\n".join(content_parts)
            
            else:
                return f"Неподдерживаемый формат документа: {ext}. Поддерживаются: DOCX, XLSX, PPTX"
                
        except Exception as e:
            logger.error(f"Error processing document {filename}: {e}")
            return f"Критическая ошибка при обработке документа {filename}: {str(e)}"

    # Обрабатываем файлы (документы, код, текст)
    if files:
        text_files_content = []
        
        for upload_file in files:
            try:
                content = await upload_file.read()
                filename = upload_file.filename or "unknown_file"
                mime_type = determine_file_mime_type(upload_file)
                
                logger.info(f"Processing file: {filename}, content_type: {upload_file.content_type}, determined_mime: {mime_type}, size: {len(content)}")
                
                # Если это PDF - извлекаем текст через Agno PDFReader
                if mime_type == "application/pdf":
                    if not AGNO_PDF_AVAILABLE:
                        # Fallback: создаем File объект для OpenAI
                        try:
                            agno_file = File(content=content, mime_type=mime_type)
                            processed_files.append(agno_file)
                            logger.info(f"Processed PDF file via File object: {filename}, size: {len(content)} bytes")
                        except Exception as fallback_error:
                            logger.error(f"Failed to process PDF {filename} as File object: {fallback_error}")
                    else:
                        # Используем Agno PDFReader для извлечения текста
                        try:
                            import io
                            pdf_reader = PDFReader()
                            documents = pdf_reader.read(io.BytesIO(content))
                            
                            if documents:
                                # Объединяем текст со всех страниц
                                pdf_text_parts = []
                                for doc in documents:
                                    if doc.content and doc.content.strip():
                                        page_num = doc.meta_data.get('page', 'неизвестно')
                                        pdf_text_parts.append(f"Страница {page_num}:\n{doc.content.strip()}")
                                
                                if pdf_text_parts:
                                    pdf_full_text = "\n\n".join(pdf_text_parts)
                                    text_files_content.append(f"""
📄 **PDF документ: {filename}**
```
{pdf_full_text}
```
""")
                                    logger.info(f"Extracted text from PDF: {filename}, {len(documents)} pages, {len(pdf_full_text)} chars")
                                else:
                                    text_files_content.append(f"📄 **PDF документ: {filename}** - Текстовое содержимое не найдено")
                                    logger.warning(f"PDF {filename} contains no extractable text")
                            else:
                                text_files_content.append(f"📄 **PDF документ: {filename}** - Не удалось извлечь содержимое")
                                logger.warning(f"Failed to extract any documents from PDF {filename}")
                                
                        except Exception as pdf_error:
                            logger.error(f"Error extracting text from PDF {filename}: {pdf_error}")
                            # Fallback: создаем File объект для OpenAI
                            try:
                                agno_file = File(content=content, mime_type=mime_type)
                                processed_files.append(agno_file)
                                logger.info(f"PDF text extraction failed, using File object: {filename}")
                            except Exception as fallback_error:
                                logger.error(f"Failed to process PDF {filename} both ways: {fallback_error}")
                                text_files_content.append(f"📄 **PDF документ: {filename}** - Ошибка обработки: {pdf_error}")
                
                # Если это документ (DOCX, XLSX, PPTX) - обрабатываем через специальные readers
                elif is_document_file(filename):
                    try:
                        document_content = await process_document_file(content, filename)
                        text_files_content.append(f"""
📄 **Документ: {filename}**
```
{document_content}
```
""")
                        logger.info(f"Processed document file: {filename}, size: {len(content)} bytes")
                    except Exception as doc_error:
                        logger.warning(f"Could not process document {filename}: {doc_error}")
                        text_files_content.append(f"📄 **Документ: {filename}** - Ошибка обработки: {doc_error}")
                
                # Если это текстовый файл - добавляем содержимое в сообщение
                elif is_text_file(filename, upload_file.content_type or ""):
                    try:
                        # Пытаемся декодировать как текст
                        text_content = content.decode('utf-8', errors='replace')
                        text_files_content.append(f"""
📄 **Файл: {filename}**
```
{text_content}
```
""")
                        logger.info(f"Processed text file: {filename}, size: {len(content)} bytes")
                    except Exception as decode_error:
                        logger.warning(f"Could not decode {filename} as text: {decode_error}")
                        text_files_content.append(f"📄 **Файл: {filename}** - Не удалось прочитать содержимое")
                
                else:
                    # Неподдерживаемый тип файла
                    logger.warning(f"Unsupported file type: {filename} ({upload_file.content_type})")
                    text_files_content.append(f"📄 **Файл: {filename}** - Неподдерживаемый тип файла")
                
            except Exception as e:
                logger.error(f"Error processing file {upload_file.filename}: {e}")
                continue
        
        # Объединяем содержимое текстовых файлов
        if text_files_content:
            text_content_summary = "\n".join(text_files_content)

    # Обрабатываем изображения
    if images:
        for upload_file in images:
            try:
                content = await upload_file.read()
                
                # Определяем формат по расширению
                format_type = None
                if upload_file.filename and '.' in upload_file.filename:
                    format_type = upload_file.filename.split('.')[-1].lower()
                
                # Создаем Agno Image объект
                agno_image = Image(
                    content=content,  # Agno ожидает raw bytes для изображений
                    format=format_type,
                    detail="auto"  # Для максимального качества анализа
                )
                processed_images.append(agno_image)
                
                logger.info(f"Processed image: {upload_file.filename}, size: {len(content)} bytes, format: {format_type}")
                
            except Exception as e:
                logger.error(f"Error processing image {upload_file.filename}: {e}")
                continue

    # Обрабатываем аудио файлы с транскрипцией
    if audio:
        for upload_file in audio:
            try:
                content = await upload_file.read()
                
                # Определяем формат по расширению
                format_type = None
                if upload_file.filename and '.' in upload_file.filename:
                    format_type = upload_file.filename.split('.')[-1].lower()
                
                # 🆕 Транскрибируем аудио через Whisper
                transcript = await transcribe_audio_file(content, upload_file.filename or "audio.mp3")
                
                if transcript:
                    # Добавляем транскрипцию к текстовому содержимому
                    text_content_summary += f"\n\n--- Транскрипция аудио файла {upload_file.filename} ---\n{transcript}\n"
                    logger.info(f"Transcribed audio: {upload_file.filename}, transcript length: {len(transcript)} chars")
                else:
                    logger.warning(f"Failed to transcribe audio: {upload_file.filename}")
                
                # Все равно создаем Audio объект (хотя он может не использоваться в Responses API)
                agno_audio = Audio(
                    content=content,  # Agno ожидает raw bytes
                    format=format_type
                )
                processed_audio.append(agno_audio)
                
                logger.info(f"Processed audio: {upload_file.filename}, size: {len(content)} bytes, format: {format_type}")
                
            except Exception as e:
                logger.error(f"Error processing audio {upload_file.filename}: {e}")
                continue

    # Обрабатываем видео файлы
    if videos:
        for upload_file in videos:
            try:
                content = await upload_file.read()
                
                # Определяем формат по расширению
                format_type = None
                if upload_file.filename and '.' in upload_file.filename:
                    format_type = upload_file.filename.split('.')[-1].lower()
                
                # Создаем Agno Video объект
                agno_video = Video(
                    content=content,  # Agno ожидает raw bytes
                    format=format_type
                )
                processed_videos.append(agno_video)
                
                logger.info(f"Processed video: {upload_file.filename}, size: {len(content)} bytes, format: {format_type}")
                
            except Exception as e:
                logger.error(f"Error processing video {upload_file.filename}: {e}")
                continue

    return (
        processed_files if processed_files else None,
        processed_images if processed_images else None,
        processed_audio if processed_audio else None,
        processed_videos if processed_videos else None,
        text_content_summary
    )


class MediaFile(BaseModel):
    """Модель для передачи файлов через JSON"""
    filename: str = Field(..., description="Имя файла")
    content: str = Field(..., description="Содержимое файла в base64")
    mime_type: Optional[str] = Field(None, description="MIME тип файла")


class RunRequest(BaseModel):
    """Request model for running an agent"""
    message: str
    stream: bool = False
    model: Model = Model.gpt_4_1
    user_id: Optional[str] = None
    session_id: Optional[str] = None
    # 🆕 Поддержка файлов через JSON
    files: Optional[List[MediaFile]] = Field(None, description="Файлы для обработки (PDF, DOCX, XLSX, PPTX, TXT, CSV, JSON, RTF)")
    images: Optional[List[MediaFile]] = Field(None, description="Изображения для анализа (PNG, JPEG, JPG, WebP)")
    audio: Optional[List[MediaFile]] = Field(None, description="Аудио файлы (MP3, WAV) - автоматическая транскрипция")
    videos: Optional[List[MediaFile]] = Field(None, description="Видео файлы (MP4, MOV, AVI) - экспериментальная поддержка")


async def convert_media_files_to_agno_objects(
    files: Optional[List[MediaFile]] = None,
    images: Optional[List[MediaFile]] = None,
    audio: Optional[List[MediaFile]] = None,
    videos: Optional[List[MediaFile]] = None
) -> tuple[Optional[List[File]], Optional[List[Image]], Optional[List[Audio]], Optional[List[Video]], str]:
    """
    Конвертирует MediaFile объекты в нативные Agno медиа объекты.
    Максимально нативное решение с Agno.
    """
    agno_files = []
    agno_images = []
    agno_audio = []
    agno_videos = []
    text_content = ""
    
    try:
        # Обрабатываем файлы
        if files:
            for media_file in files:
                try:
                    # Декодируем base64 содержимое
                    file_content = base64.b64decode(media_file.content)
                    
                    # Проверяем тип файла
                    if media_file.mime_type == "application/pdf":
                        # PDF файлы - извлекаем текст через Agno PDFReader
                        if not AGNO_PDF_AVAILABLE:
                            # Fallback: создаем File объект для OpenAI
                            try:
                                agno_file = File(content=file_content, mime_type=media_file.mime_type)
                                agno_files.append(agno_file)
                                logger.info(f"Processed PDF file via File object: {media_file.filename}, size: {len(file_content)} bytes")
                            except Exception as fallback_error:
                                logger.error(f"Failed to process PDF {media_file.filename} as File object: {fallback_error}")
                        else:
                            # Используем Agno PDFReader для извлечения текста
                            try:
                                import io
                                pdf_reader = PDFReader()
                                documents = pdf_reader.read(io.BytesIO(file_content))
                                
                                if documents:
                                    # Объединяем текст со всех страниц
                                    pdf_text_parts = []
                                    for doc in documents:
                                        if doc.content and doc.content.strip():
                                            page_num = doc.meta_data.get('page', 'неизвестно')
                                            pdf_text_parts.append(f"Страница {page_num}:\n{doc.content.strip()}")
                                    
                                    if pdf_text_parts:
                                        pdf_full_text = "\n\n".join(pdf_text_parts)
                                        text_content += f"\n\n--- Содержимое PDF документа {media_file.filename} ---\n{pdf_full_text}\n"
                                        logger.info(f"Extracted text from PDF: {media_file.filename}, {len(documents)} pages, {len(pdf_full_text)} chars")
                                    else:
                                        text_content += f"\n\n--- PDF документ {media_file.filename} ---\nТекстовое содержимое не найдено\n"
                                        logger.warning(f"PDF {media_file.filename} contains no extractable text")
                                else:
                                    text_content += f"\n\n--- PDF документ {media_file.filename} ---\nНе удалось извлечь содержимое\n"
                                    logger.warning(f"Failed to extract any documents from PDF {media_file.filename}")
                                    
                            except Exception as pdf_error:
                                logger.error(f"Error extracting text from PDF {media_file.filename}: {pdf_error}")
                                # Fallback: создаем File объект для OpenAI
                                try:
                                    agno_file = File(content=file_content, mime_type=media_file.mime_type)
                                    agno_files.append(agno_file)
                                    logger.info(f"PDF text extraction failed, using File object: {media_file.filename}")
                                except Exception as fallback_error:
                                    logger.error(f"Failed to process PDF {media_file.filename} both ways: {fallback_error}")
                                    text_content += f"\n\n--- PDF документ {media_file.filename} ---\nОшибка обработки: {pdf_error}\n"
                    
                    elif media_file.mime_type in [
                        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",  # DOCX
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",       # XLSX
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation" # PPTX
                    ]:
                        # Документы Office обрабатываем через специальные readers
                        try:
                            # Определяем расширение по MIME типу
                            if "wordprocessing" in media_file.mime_type:
                                ext = "docx"
                            elif "spreadsheet" in media_file.mime_type:
                                ext = "xlsx"
                            elif "presentation" in media_file.mime_type:
                                ext = "pptx"
                            else:
                                ext = "unknown"
                            
                            # Создаем временное имя файла с правильным расширением
                            temp_filename = f"{media_file.filename}.{ext}" if not media_file.filename.endswith(f".{ext}") else media_file.filename
                            
                            # Обрабатываем документ
                            from api.routes.agents import process_document_file
                            document_content = await process_document_file(file_content, temp_filename)
                            text_content += f"\n\n--- Содержимое документа {media_file.filename} ---\n{document_content}\n"
                            logger.info(f"Processed document file: {media_file.filename}, size: {len(file_content)} bytes")
                            
                        except Exception as doc_error:
                            logger.warning(f"Could not process document {media_file.filename}: {doc_error}")
                            text_content += f"\n\n--- Документ {media_file.filename} ---\nОшибка обработки: {doc_error}\n"
                    
                    elif media_file.mime_type and (
                        media_file.mime_type.startswith('text/') or
                        media_file.mime_type in ['application/json', 'application/csv', 'text/rtf']
                    ):
                        # Текстовые файлы извлекаем содержимое
                        try:
                            file_text = file_content.decode('utf-8')
                            text_content += f"\n\n--- Содержимое файла {media_file.filename} ---\n{file_text}\n"
                            logger.info(f"Processed text file: {media_file.filename}, size: {len(file_content)} bytes")
                        except UnicodeDecodeError:
                            logger.warning(f"Could not decode text content from {media_file.filename}")
                    
                    else:
                        # Остальные файлы пытаемся обработать как File объекты
                        try:
                            agno_file = File(
                                content=file_content,
                                mime_type=media_file.mime_type
                            )
                            agno_files.append(agno_file)
                            logger.info(f"Processed generic file: {media_file.filename}, size: {len(file_content)} bytes")
                        except Exception as file_error:
                            logger.warning(f"Could not process file {media_file.filename}: {file_error}")
                            text_content += f"\n\n--- Файл {media_file.filename} ---\nНе удалось обработать: {file_error}\n"
                            
                except Exception as e:
                    logger.error(f"Error processing file {media_file.filename}: {e}")
                    continue
        
        # Обрабатываем изображения
        if images:
            for media_file in images:
                try:
                    image_content = base64.b64decode(media_file.content)
                    agno_image = Image(
                        content=image_content,
                        name=media_file.filename,
                        mime_type=media_file.mime_type
                    )
                    agno_images.append(agno_image)
                except Exception as e:
                    logger.error(f"Error processing image {media_file.filename}: {e}")
                    continue
        
        # Обрабатываем аудио с транскрипцией
        if audio:
            for media_file in audio:
                try:
                    audio_content = base64.b64decode(media_file.content)
                    
                    # 🆕 Транскрибируем аудио через Whisper
                    transcript = await transcribe_audio_file(audio_content, media_file.filename)
                    
                    if transcript:
                        # Добавляем транскрипцию к текстовому содержимому
                        text_content += f"\n\n--- Транскрипция аудио файла {media_file.filename} ---\n{transcript}\n"
                        logger.info(f"Transcribed audio: {media_file.filename}, transcript length: {len(transcript)} chars")
                    else:
                        logger.warning(f"Failed to transcribe audio: {media_file.filename}")
                    
                    # Все равно создаем Audio объект
                    agno_audio_obj = Audio(
                        content=audio_content,
                        name=media_file.filename,
                        mime_type=media_file.mime_type
                    )
                    agno_audio.append(agno_audio_obj)
                except Exception as e:
                    logger.error(f"Error processing audio {media_file.filename}: {e}")
                    continue
        
        # Обрабатываем видео
        if videos:
            for media_file in videos:
                try:
                    video_content = base64.b64decode(media_file.content)
                    agno_video = Video(
                        content=video_content,
                        name=media_file.filename,
                        mime_type=media_file.mime_type
                    )
                    agno_videos.append(agno_video)
                except Exception as e:
                    logger.error(f"Error processing video {media_file.filename}: {e}")
                    continue
        
    except Exception as e:
        logger.error(f"Error converting media files: {e}")
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Error processing media files: {str(e)}"
        )
    
    return (
        agno_files if agno_files else None,
        agno_images if agno_images else None,
        agno_audio if agno_audio else None,
        agno_videos if agno_videos else None,
        text_content.strip()
    )


async def _execute_agent_run(
    agent_id: str,
    message: str,
    stream: bool = False,
    model: str = "gpt-4.1",
    user_id: Optional[str] = None,
    session_id: Optional[str] = None,
    files: Optional[Union[List[MediaFile], List[UploadFile]]] = None,
    images: Optional[Union[List[MediaFile], List[UploadFile]]] = None,
    audio: Optional[Union[List[MediaFile], List[UploadFile]]] = None,
    videos: Optional[Union[List[MediaFile], List[UploadFile]]] = None,
):
    """
    Универсальная функция выполнения запроса к агенту.
    Поддерживает как MediaFile (JSON), так и UploadFile (multipart).
    """
    try:
        # Получаем агента
        agent = get_agent(agent_id)
        if not agent:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail=f"Agent '{agent_id}' not found"
            )

        # Определяем тип файлов и обрабатываем соответственно
        has_any_media = any([
            files and len(files) > 0,
            images and len(images) > 0,
            audio and len(audio) > 0, 
            videos and len(videos) > 0
        ])
        
        if has_any_media:
            # Находим первый непустой список для определения типа
            first_media_item = None
            for media_list in [files, images, audio, videos]:
                if media_list and len(media_list) > 0:
                    first_media_item = media_list[0]
                    break
            
            if first_media_item and isinstance(first_media_item, MediaFile):
                # JSON формат с base64
                logger.info(f"Processing media as JSON/base64 format")
                processed_files, processed_images, processed_audio, processed_videos, text_content = await convert_media_files_to_agno_objects(
                    files=files, images=images, audio=audio, videos=videos
                )
            else:
                # Multipart формат с UploadFile
                logger.info(f"Processing media as multipart/UploadFile format")
                processed_files, processed_images, processed_audio, processed_videos, text_content = await process_uploaded_files(
                    files=files, images=images, audio=audio, videos=videos
                )
        else:
            # Нет файлов
            logger.info(f"No media files to process")
            processed_files = processed_images = processed_audio = processed_videos = None
            text_content = ""

        # Объединяем текстовое сообщение с содержимым файлов
        enhanced_message = message
        if text_content:
            enhanced_message = f"{message}\n\n{text_content}"

        # Настройка модели
        if hasattr(agent, 'model') and agent.model:
            if hasattr(agent.model, 'id'):
                agent.model.id = model

        # Запускаем агента
        if stream:
            return StreamingResponse(
                chat_response_streamer(
                    agent, enhanced_message, 
                    files=processed_files, images=processed_images,
                    audio=processed_audio, videos=processed_videos
                ),
                media_type="text/event-stream",
            )
        else:
            response = await agent.arun(
                enhanced_message, stream=False,
                files=processed_files, images=processed_images,
                audio=processed_audio, videos=processed_videos
            )
            
            # 🆕 ПОЛНАЯ ПОДДЕРЖКА ВЫХОДНЫХ МУЛЬТИМЕДИА
            def safe_extract_multimedia_content(content):
                """Безопасно извлекает содержимое, обрабатывая bytes и строки"""
                if content is None:
                    return None
                if isinstance(content, bytes):
                    try:
                        return content.decode('utf-8')
                    except UnicodeDecodeError:
                        import base64
                        return base64.b64encode(content).decode('ascii')
                return content
            
            return {
                "content": getattr(response, 'content', ''),
                "content_type": getattr(response, 'content_type', 'text/plain'),
                
                # 🎨 Сгенерированные изображения
                "images": [
                    {
                        "id": getattr(img, 'id', None),
                        "url": getattr(img, 'url', None),
                        "content": safe_extract_multimedia_content(getattr(img, 'content', None)),
                        "mime_type": getattr(img, 'mime_type', None),
                        "alt_text": getattr(img, 'alt_text', None)
                    } for img in getattr(response, 'images', [])
                ] if hasattr(response, 'images') and response.images else None,
                
                # 🎵 Сгенерированные аудио файлы
                "audio": [
                    {
                        "id": getattr(aud, 'id', None),
                        "url": getattr(aud, 'url', None),
                        "content": getattr(aud, 'base64_audio', None),
                        "mime_type": getattr(aud, 'mime_type', None),
                        "length": getattr(aud, 'length', None)
                    } for aud in getattr(response, 'audio', [])
                ] if hasattr(response, 'audio') and response.audio else None,
                
                # 🎬 Сгенерированные видео файлы
                "videos": [
                    {
                        "id": getattr(vid, 'id', None),
                        "url": getattr(vid, 'url', None),
                        "content": safe_extract_multimedia_content(getattr(vid, 'content', None)),
                        "mime_type": getattr(vid, 'mime_type', None),
                        "eta": getattr(vid, 'eta', None),
                        "length": getattr(vid, 'length', None)
                    } for vid in getattr(response, 'videos', [])
                ] if hasattr(response, 'videos') and response.videos else None,
                
                # 🎤 Голосовой ответ агента (TTS)
                "response_audio": {
                    "id": getattr(response.response_audio, 'id', None),
                    "content": getattr(response.response_audio, 'content', None),
                    "expires_at": getattr(response.response_audio, 'expires_at', None),
                    "transcript": getattr(response.response_audio, 'transcript', None),
                    "mime_type": getattr(response.response_audio, 'mime_type', None),
                    "sample_rate": getattr(response.response_audio, 'sample_rate', None),
                    "channels": getattr(response.response_audio, 'channels', None)
                } if hasattr(response, 'response_audio') and response.response_audio else None,
                
                # 📚 Цитаты и источники
                "citations": response.citations.model_dump() if hasattr(response, 'citations') and response.citations else None,
                
                # 🧠 Мышление и рассуждения
                "thinking": getattr(response, 'thinking', None),
                "reasoning_content": getattr(response, 'reasoning_content', None),
                
                # 📊 Метрики и информация
                "metrics": getattr(response, 'metrics', None),
                "model": getattr(response, 'model', None),
                "run_id": getattr(response, 'run_id', None),
                "session_id": getattr(response, 'session_id', None),
                "formatted_tool_calls": getattr(response, 'formatted_tool_calls', None)
            }

    except Exception as e:
        logger.error(f"Error running agent {agent_id}: {e}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Agent execution failed: {str(e)}"
        )


@agents_router.post("/{agent_id}/runs", status_code=status.HTTP_200_OK)
async def create_agent_run_universal(
    agent_id: str,
    request: Request,
    # JSON параметры (используются если Content-Type: application/json)
    body: Optional[RunRequest] = None,
    # Multipart параметры (используются если Content-Type: multipart/form-data)
    message: Optional[str] = Form(None, description="Текстовое сообщение для агента"),
    stream: Optional[bool] = Form(False, description="Включить потоковый ответ"),
    model: Optional[str] = Form("gpt-4.1", description="Модель для использования"),
    user_id: Optional[str] = Form(None, description="ID пользователя"),
    session_id: Optional[str] = Form(None, description="ID сессии"),
    files: Optional[List[UploadFile]] = FastAPIFile(None, description="Файлы для обработки (PDF, DOCX, XLSX, PPTX, TXT, CSV, JSON, RTF)"),
    images: Optional[List[UploadFile]] = FastAPIFile(None, description="Изображения для анализа (PNG, JPEG, JPG, WebP)"),
    audio: Optional[List[UploadFile]] = FastAPIFile(None, description="Аудио файлы (MP3, WAV) - автоматическая транскрипция"),
    videos: Optional[List[UploadFile]] = FastAPIFile(None, description="Видео файлы (MP4, MOV, AVI) - экспериментальная поддержка"),
):
    """
    🚀 УНИВЕРСАЛЬНЫЙ ЭНДПОИНТ для запуска агентов.
    
    Автоматически определяет формат запроса:
    - Content-Type: application/json → JSON с base64 файлами
    - Content-Type: multipart/form-data → Файлы + форма
    
    Примеры использования:
    
    1. JSON формат (удобно для веб-приложений):
    ```json
    {
        "message": "Проанализируй документ",
        "files": [{"filename": "doc.pdf", "content": "JVBERi0x...", "mime_type": "application/pdf"}]
    }
    ```
    
    2. Multipart формат (эффективно для больших файлов):
    ```bash
    curl -X POST "/v1/agents/finance/runs" \
         -F "message=Проанализируй документ" \
         -F "files=@document.pdf"
    ```
    """
    content_type = request.headers.get("content-type", "")
    
    if content_type.startswith("multipart/form-data"):
        # Multipart запрос
        logger.info(f"🚀 Multipart request: agent={agent_id}, message_len={len(message) if message else 0}, "
                    f"files={len(files) if files else 0}, images={len(images) if images else 0}, "
                    f"audio={len(audio) if audio else 0}, videos={len(videos) if videos else 0}")
        
        return await _execute_agent_run(
            agent_id=agent_id,
            message=message,
            stream=stream,
            model=model,
            user_id=user_id,
            session_id=session_id,
            files=files,
            images=images,
            audio=audio,
            videos=videos
        )
    else:
        # JSON запрос
        if body is None:
            # Если body не передан, пытаемся распарсить из request
            try:
                request_data = await request.json()
                body = RunRequest(**request_data)
            except Exception as e:
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail=f"Invalid JSON request: {str(e)}"
                )
        
        logger.info(f"📄 JSON request: agent={agent_id}, message_len={len(body.message)}, "
                    f"files={len(body.files) if body.files else 0}, "
                    f"images={len(body.images) if body.images else 0}, "
                    f"audio={len(body.audio) if body.audio else 0}, "
                    f"videos={len(body.videos) if body.videos else 0}")
        
        return await _execute_agent_run(
            agent_id=agent_id,
            message=body.message,
            stream=body.stream,
            model=body.model.value,
            user_id=body.user_id,
            session_id=body.session_id,
            files=body.files,
            images=body.images,
            audio=body.audio,
            videos=body.videos
        )


@agents_router.get("/{agent_id}/sessions", status_code=status.HTTP_200_OK)
async def get_agent_sessions(agent_id: str, user_id: Optional[str] = None):
    """
    Получение сессий для конкретного агента.
    Аналогично native Agno endpoint AGENT_SESSION_CREATE.

    Args:
        agent_id: ID агента для получения сессий
        user_id: Опциональный ID пользователя для фильтрации

    Returns:
        Список сессий агента с метаданными
    """
    try:
        # Проверяем существование агента
        agent: Agent = get_agent(agent_id=agent_id)
        if not agent:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail=f"Agent {agent_id} not found"
            )

        # Получаем сессии из storage если он настроен
        sessions = []
        if hasattr(agent, 'storage') and agent.storage:
            try:
                # Пытаемся получить сессии из storage
                # В Agno storage обычно содержит методы для работы с сессиями
                if hasattr(agent.storage, 'get_sessions'):
                    sessions_data = agent.storage.get_sessions(
                        agent_id=agent_id, 
                        user_id=user_id
                    )
                    if sessions_data:
                        sessions = [
                            {
                                "session_id": session.session_id,
                                "user_id": session.user_id,
                                "agent_id": session.agent_id,
                                "created_at": session.created_at,
                                "updated_at": session.updated_at,
                                "session_data": session.session_data,
                                "agent_data": session.agent_data
                            }
                            for session in sessions_data
                            if isinstance(session, AgentSession)
                        ]
            except Exception as e:
                logger.warning(f"Could not retrieve sessions from storage: {e}")
                # Если не удается получить из storage, возвращаем пустой список

        # Возвращаем ответ в формате совместимом с Agno
        return {
            "agent_id": agent_id,
            "user_id": user_id,
            "sessions": sessions,
            "total_sessions": len(sessions),
            "agent_type": "static" if hasattr(agent, 'agent_id') else "dynamic",
            "timestamp": datetime.utcnow().isoformat()
        }

    except HTTPException:
        raise
    except ValueError as e:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail=str(e))
    except Exception as e:
        logger.error(f"Error retrieving sessions for agent {agent_id}: {e}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to retrieve sessions for agent {agent_id}"
        )


@agents_router.post("/{agent_id}/knowledge/load", status_code=status.HTTP_200_OK)
async def load_agent_knowledge(agent_id: str):
    """
    Loads the knowledge base for a specific agent.

    Args:
        agent_id: The ID of the agent to load knowledge for.

    Returns:
        A success message if the knowledge base is loaded.
    """
    agent_knowledge: Optional[AgentKnowledge] = None

    if agent_id == "agno_assist":
        agent_knowledge = get_agno_assist_knowledge()
    else:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Agent {agent_id} does not have a knowledge base.",
        )

    try:
        await agent_knowledge.aload(upsert=True)
    except Exception as e:
        logger.error(f"Error loading knowledge base for {agent_id}: {e}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to load knowledge base for {agent_id}.",
        )

    return {"message": f"Knowledge base for {agent_id} loaded successfully."}


# 🆕 Функция для транскрипции аудио через Whisper
async def transcribe_audio_file(audio_content: bytes, filename: str) -> Optional[str]:
    """
    Транскрибирует аудио файл через OpenAI Whisper API
    
    Args:
        audio_content: Содержимое аудио файла в байтах
        filename: Имя файла для определения формата
        
    Returns:
        Текст транскрипции или None в случае ошибки
    """
    try:
        # Создаем временный файл
        with tempfile.NamedTemporaryFile(delete=False, suffix=Path(filename).suffix) as temp_file:
            temp_file.write(audio_content)
            temp_file_path = temp_file.name
        
        try:
            # Инициализируем OpenAI клиент для Whisper
            client = AsyncOpenAI(api_key=os.getenv("OPENAI_API_KEY"))
            
            # Открываем файл для Whisper API
            with open(temp_file_path, "rb") as audio_file:
                transcript = await client.audio.transcriptions.create(
                    model="whisper-1",
                    file=audio_file,
                    language="ru"  # Можно сделать автоопределение или параметром
                )
            
            return transcript.text
            
        finally:
            # Удаляем временный файл
            os.unlink(temp_file_path)
            
    except Exception as e:
        print(f"Ошибка при транскрипции аудио: {e}")
        return None
